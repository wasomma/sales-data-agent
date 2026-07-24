"""Ingest Excel and PowerPoint files from the inbox into DuckDB.

Files are snapshots: every row carries as_of_date + source_file + file_hash, and
re-running sync on unchanged files is a no-op (hash match in ingest_log).
"""

import hashlib
import json
import re
from datetime import date, datetime
from pathlib import Path

import pandas as pd
import yaml
from pptx import Presentation

from .config import INBOX_DIR, MAPPING_PATH
from .db import connect_rw

DATE_IN_NAME = re.compile(r"(\d{4}-\d{2}-\d{2})")

CANONICAL_COLUMNS = [
    "deal_name", "account", "amount", "close_date",
    "stage", "owner", "forecast_category",
]
# A sheet/table must map at least these to be treated as deal data.
REQUIRED_COLUMNS = {"deal_name", "account", "amount"}


def load_mapping() -> dict[str, str]:
    """Return {lowercased source header -> canonical name}."""
    raw = yaml.safe_load(MAPPING_PATH.read_text(encoding="utf-8"))
    lookup: dict[str, str] = {}
    for canonical, aliases in raw["column_aliases"].items():
        lookup[canonical.lower()] = canonical
        for alias in aliases:
            lookup[alias.strip().lower()] = canonical
    return lookup


def file_hash(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()[:16]


def infer_as_of_date(path: Path) -> date:
    """Priority: YYYY-MM-DD in filename, else file modification time."""
    m = DATE_IN_NAME.search(path.name)
    if m:
        return date.fromisoformat(m.group(1))
    return datetime.fromtimestamp(path.stat().st_mtime).date()


def deal_key(account: str, deal_name: str) -> str:
    # Stable identity across snapshots. Falls back to name-based hash because
    # exports carry no CRM id (documented limitation in DESIGN.md).
    basis = f"{account.strip().lower()}|{deal_name.strip().lower()}"
    return hashlib.md5(basis.encode()).hexdigest()[:12]


def map_columns(df: pd.DataFrame, lookup: dict[str, str]) -> pd.DataFrame | None:
    """Rename recognized columns to canonical names; None if not a deals table."""
    renames = {}
    for col in df.columns:
        canonical = lookup.get(str(col).strip().lower())
        if canonical:
            renames[col] = canonical
    if not REQUIRED_COLUMNS.issubset(renames.values()):
        return None
    df = df.rename(columns=renames)
    return df[[c for c in CANONICAL_COLUMNS if c in df.columns]]


def normalize_rows(df: pd.DataFrame) -> tuple[list[dict], list[tuple[dict, str]]]:
    """Validate and coerce rows. Returns (good_rows, [(raw_row, reason), ...])."""
    good: list[dict] = []
    rejects: list[tuple[dict, str]] = []
    for _, row in df.iterrows():
        raw = {k: (None if pd.isna(v) else v) for k, v in row.items()}
        name, account = raw.get("deal_name"), raw.get("account")
        if not name or not account:
            rejects.append((raw, "missing deal_name or account"))
            continue
        try:
            amount = round(float(str(raw["amount"]).replace("$", "").replace(",", "")), 2)
        except (TypeError, ValueError):
            rejects.append((raw, f"unparseable amount: {raw.get('amount')!r}"))
            continue
        close = raw.get("close_date")
        if close is not None:
            try:
                close = pd.to_datetime(close).date()
            except (ValueError, TypeError):
                rejects.append((raw, f"unparseable close_date: {raw.get('close_date')!r}"))
                continue
        good.append({
            "deal_key": deal_key(str(account), str(name)),
            "deal_name": str(name).strip(),
            "account": str(account).strip(),
            "amount": amount,
            "close_date": close,
            "stage": str(raw["stage"]).strip() if raw.get("stage") else None,
            "owner": str(raw["owner"]).strip() if raw.get("owner") else None,
            "forecast_category": str(raw["forecast_category"]).strip()
            if raw.get("forecast_category") else None,
        })
    return good, rejects


def frames_from_xlsx(path: Path) -> list[pd.DataFrame]:
    sheets = pd.read_excel(path, sheet_name=None)
    return list(sheets.values())


def frames_from_pptx(path: Path) -> tuple[list[pd.DataFrame], list[tuple[int, str]]]:
    """Extract tables as DataFrames and collect narrative slide text."""
    prs = Presentation(str(path))
    frames: list[pd.DataFrame] = []
    texts: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in r.cells] for r in shape.table.rows]
                if len(rows) >= 2:
                    frames.append(pd.DataFrame(rows[1:], columns=rows[0]))
            elif shape.has_text_frame and shape.text_frame.text.strip():
                chunks.append(shape.text_frame.text.strip())
        if chunks:
            texts.append((i, "\n".join(chunks)))
    return frames, texts


def sync() -> list[dict]:
    """Ingest new/changed inbox files. Returns a per-file result summary."""
    lookup = load_mapping()
    results = []
    con = connect_rw()
    try:
        known = {h for (h,) in con.execute("SELECT file_hash FROM ingest_log").fetchall()}
        files = sorted(p for p in INBOX_DIR.glob("*")
                       if p.suffix.lower() in (".xlsx", ".pptx") and not p.name.startswith("~"))
        for path in files:
            fhash = file_hash(path)
            if fhash in known:
                results.append({"file": path.name, "status": "unchanged", "rows": 0, "rejected": 0})
                continue
            as_of = infer_as_of_date(path)

            if path.suffix.lower() == ".xlsx":
                frames, texts = frames_from_xlsx(path), []
            else:
                frames, texts = frames_from_pptx(path)

            loaded = rejected = 0
            for df in frames:
                mapped = map_columns(df, lookup)
                if mapped is None:
                    continue
                good, bad = normalize_rows(mapped)
                for r in good:
                    con.execute(
                        "INSERT INTO deals VALUES (?,?,?,?,?,?,?,?,?,?,?)",
                        [r["deal_key"], r["deal_name"], r["account"], r["amount"],
                         r["close_date"], r["stage"], r["owner"], r["forecast_category"],
                         as_of, path.name, fhash],
                    )
                for raw, reason in bad:
                    con.execute("INSERT INTO rejects VALUES (?,?,?)",
                                [path.name, json.dumps(raw, default=str), reason])
                loaded += len(good)
                rejected += len(bad)

            for slide_no, text in texts:
                con.execute("INSERT INTO slide_text VALUES (?,?,?,?)",
                            [path.name, slide_no, text, as_of])

            con.execute(
                "INSERT INTO ingest_log (source_file, file_hash, as_of_date, rows_loaded, rows_rejected)"
                " VALUES (?,?,?,?,?)",
                [path.name, fhash, as_of, loaded, rejected],
            )
            results.append({"file": path.name, "status": "ingested",
                            "rows": loaded, "rejected": rejected, "as_of": str(as_of)})
    finally:
        con.close()
    return results
