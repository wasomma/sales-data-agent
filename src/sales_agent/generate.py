"""Deterministic synthetic sales data generator.

Produces realistic CRM-export-shaped files in the inbox:
  - two Excel pipeline snapshots (different as_of dates, with drift between them
    to exercise snapshot logic: amounts change, stages advance, deals slip/close)
  - one QBR PowerPoint with an embedded top-deals table and narrative text

Seeded RNG -> same dataset every run, so answers are reproducible in tests/demos.
"""

import random
from datetime import date, timedelta
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from .config import INBOX_DIR, ensure_dirs

SEED = 20260724

ACCOUNTS = [
    "Northwind Logistics", "Apex Manufacturing", "BlueSky Airlines", "Cascade Health",
    "Delta Retail Group", "Evergreen Energy", "Fairbanks Financial", "Granite Insurance",
    "Harborview Hotels", "Ironclad Security", "Juniper Media", "Keystone Pharma",
]
OWNERS = ["Alex Chen", "Priya Patel", "Marcus Webb", "Sofia Reyes", "Dan Kowalski"]
STAGES = ["Prospecting", "Qualification", "Proposal", "Negotiation", "Closed Won", "Closed Lost"]
PRODUCTS = ["Platform License", "Analytics Add-on", "Enterprise Rollout", "Renewal",
            "Pilot Program", "Support Upgrade", "Data Migration", "Expansion"]

SNAPSHOT_1 = date(2026, 6, 30)
SNAPSHOT_2 = date(2026, 7, 21)


def _build_deals(rng: random.Random, n: int = 120) -> list[dict]:
    deals = []
    for i in range(n):
        account = rng.choice(ACCOUNTS)
        product = rng.choice(PRODUCTS)
        name = f"{account.split()[0]} {product} {rng.randint(100, 999)}"
        close = date(2026, 1, 1) + timedelta(days=rng.randint(0, 540))  # 2026 through mid-2027
        stage = rng.choices(STAGES, weights=[20, 20, 20, 15, 15, 10])[0]
        deals.append({
            "Opportunity Name": name,
            "Account Name": account,
            "Amount": rng.choice([15, 25, 40, 60, 85, 120, 180, 250, 400, 650]) * 1000
            + rng.randint(0, 9) * 500,
            "Close Date": close,
            "Stage": stage,
            "Owner": rng.choice(OWNERS),
            "Forecast Category": (
                "Closed" if stage.startswith("Closed")
                else rng.choice(["Commit", "Best Case", "Pipeline"])
            ),
        })
    return deals


def _drift(rng: random.Random, deals: list[dict]) -> list[dict]:
    """Simulate three weeks of pipeline movement for the second snapshot."""
    stage_order = ["Prospecting", "Qualification", "Proposal", "Negotiation"]
    out = []
    for d in deals:
        d = dict(d)
        roll = rng.random()
        if d["Stage"] in stage_order:
            if roll < 0.20:  # advance a stage
                idx = stage_order.index(d["Stage"])
                d["Stage"] = (stage_order + ["Closed Won"])[idx + 1]
                if d["Stage"] == "Closed Won":
                    d["Forecast Category"] = "Closed"
            elif roll < 0.30:  # slip a quarter
                d["Close Date"] = d["Close Date"] + timedelta(days=91)
            elif roll < 0.38:  # resize
                d["Amount"] = int(d["Amount"] * rng.choice([0.8, 1.15, 1.3]))
        out.append(d)
    # a few brand-new deals appear between snapshots
    out.extend(_build_deals(rng, 8))
    return out


def _write_xlsx(deals: list[dict], path: Path) -> None:
    pd.DataFrame(deals).to_excel(path, index=False, sheet_name="Pipeline")


def _write_qbr_pptx(deals: list[dict], path: Path, as_of: date) -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Q3 2026 Pipeline Review"
    title_slide.placeholders[1].text = f"Sales Ops — data as of {as_of.isoformat()}"

    narrative = prs.slides.add_slide(blank)
    box = narrative.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tf = box.text_frame
    tf.text = "Commentary"
    p = tf.add_paragraph()
    p.text = (
        "Enterprise segment momentum continues; renewal risk concentrated in "
        "Delta Retail Group. Cascade Health expansion moved to Negotiation."
    )

    top = sorted(deals, key=lambda d: d["Amount"], reverse=True)[:10]
    table_slide = prs.slides.add_slide(blank)
    rows, cols = len(top) + 1, 5
    shape = table_slide.shapes.add_table(rows, cols, Inches(0.4), Inches(0.6),
                                         Inches(9.2), Inches(5.5))
    table = shape.table
    headers = ["Deal", "Account", "Amount", "Close Date", "Stage"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
    for r, d in enumerate(top, start=1):
        table.cell(r, 0).text = d["Opportunity Name"]
        table.cell(r, 1).text = d["Account Name"]
        table.cell(r, 2).text = str(d["Amount"])
        table.cell(r, 3).text = d["Close Date"].isoformat()
        table.cell(r, 4).text = d["Stage"]
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)

    prs.save(str(path))


def generate() -> list[str]:
    """Write the synthetic dataset into the inbox. Returns created file names."""
    ensure_dirs()
    rng = random.Random(SEED)
    snap1 = _build_deals(rng)
    snap2 = _drift(rng, snap1)

    files = {
        INBOX_DIR / f"sales_pipeline_{SNAPSHOT_1.isoformat()}.xlsx": lambda p: _write_xlsx(snap1, p),
        INBOX_DIR / f"sales_pipeline_{SNAPSHOT_2.isoformat()}.xlsx": lambda p: _write_xlsx(snap2, p),
        INBOX_DIR / f"qbr_deck_{SNAPSHOT_2.isoformat()}.pptx": lambda p: _write_qbr_pptx(snap2, p, SNAPSHOT_2),
    }
    for path, writer in files.items():
        writer(path)
    return [p.name for p in files]
